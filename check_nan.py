import pandas as pd
import requests
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import re

    
# 模糊搜索：取前5000个字符，滑动窗口为10
def find_split_code(chars, mixed_string):
    '''
    使用正则表达式匹配数字和非数字字符
    '''
    digits = re.findall(r'\d+', mixed_string)
    letters = re.findall(r'\D+', mixed_string)
    res = digits+letters
    ans = True
    for char in res:
        if char not in chars:
            ans = False
    return ans

def sliding_window(chars, window_size, course_code):
    '''
    滑动窗口寻找
    '''
    ans = False 
    if course_code in chars:
        return True
    for i in range(len(chars) - window_size + 1):
        if find_split_code(chars[i:i + window_size], course_code):
            ans = True
            break
    return ans

def check_pdf(url, course_code):
    '''
    检查PDF文件中是否出现course_code
    '''
    # 发送HTTP请求下载文件
    response = requests.get(url)
    # 检查请求是否成功
    if response.status_code == 200:
        # print('File downloaded successfully.')
        from io import BytesIO
        pdf_data = BytesIO(response.content)
        reader = PdfReader(pdf_data)
        total_chars = 0
        text_content = []
        # 遍历PDF中的每一页
        for page_num in range(len(reader.pages)):
            # 获取当前页
            page = reader.pages[page_num]
            # 获取当前页的内容
            text = page.extract_text()
            tt = []
            temp = text.split('\n')
            for item in temp:
                tt += item.split(' ')
            text_content += tt
            total_chars += len(tt)
            if total_chars >= 5000:
                break
        # 滑动窗口寻找
        return sliding_window(text_content, 10, course_code)
    else:
        # print('Failed to download the file.')
        return False
 
def check_docx(url, course_code):
    '''
    检查word文件中是否出现course_code
    '''
    # 发送HTTP请求下载文件
    response = requests.get(url)
    # 检查请求是否成功
    if response.status_code == 200:
        # print('File downloaded successfully.')
        from io import BytesIO
        docx_data = BytesIO(response.content)
        # 使用Document解析DOCX
        doc = Document(docx_data)
        total_chars = 0
        text_content = []
        # 遍历DOCX中的每一个段落
        for para in doc.paragraphs:
            tt = []
            temp = para.text.split('\n')
            for item in temp:
                tt += item.split(' ')
            text_content += tt
            total_chars += len(tt)
            if total_chars >= 5000:
                break
        # 滑动窗口寻找
        return sliding_window(text_content, 10, course_code)
    else:
        # print('Failed to download the file.')
        return False   

def check_ppt(url, course_code):
    '''
    检查PPT文件中是否出现course_code
    '''
    # 发送HTTP请求下载文件
    response = requests.get(url)
    # 检查请求是否成功
    if response.status_code == 200:
        # print('File downloaded successfully.')
        from io import BytesIO
        ppt_data = BytesIO(response.content)
        # 使用Presentation解析PPT
        prs = Presentation(ppt_data)
        total_chars = 0
        text_content = []
        # 遍历PPT中的每一个幻灯片
        for slide_num, slide in enumerate(prs.slides):
            # 遍历幻灯片中的每一个形状
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    tt = []
                    temp = shape.text.split('\n')
                    for item in temp:
                        tt += item.split(' ')
                    text_content += tt
                    total_chars += len(tt)
                    if total_chars >= 5000:
                        break
            if total_chars >= 5000:
                break
        # 滑动窗口寻找
        return sliding_window(text_content, 10, course_code)
    else:
        # print('Failed to download the file.')
        return False


def check_nan(select_df, link_df):
    '''
    主函数，对于每一个order_id，检查是否有PDF、DOCX、PPT文件中出现course_code
    '''
    select_df = select_df.reset_index(drop=True)
    select_df['test_na'] = False
    for index, row in select_df.iterrows():
        # print(index)  # 查看进度
        order_id = row['order_id']
        course_code = row['course_code']
        test = False
        temp_df = link_df[link_df['order_id']==order_id]
        for link in temp_df['oss_link']:
            if link.endswith('pdf'):
                test = check_pdf(link, str(course_code))
            elif link.endswith('docx'):
                test = check_docx(link, str(course_code))
            elif link.endswith('pptx'):
                test = check_ppt(link, str(course_code))
            if test:
                break
        select_df.loc[index, 'test_na'] = test
    return select_df


# 读取oss_link文件
link_df = pd.read_excel('syllabus_info_export_after_4may_1430.xlsx', sheet_name='Sheet3')
df = pd.read_excel('syllabus_info_export_after_4may_1430.xlsx', sheet_name='order_id')

# 算法是nan，人工不是nan，判断人工course_code是否出现在文件里
select_df1 = df[(df['couse_code_alg'].isna()) & (df['course_code'].notna())].reset_index(drop=True)
res_df1 = check_nan(select_df1, link_df)[['order_id','test_na']]
res_df1.to_excel('simp_result1_6may.xlsx', index=False)

# 人工不是nan、算法不是nan、相似度在0.5以下的course_code，查找文件中是否出现人工维护的course_code
select_df2 = df[(df['couse_code_alg'].notna()) & (df['course_code'].notna()) & (df['similarity']<0.5)].reset_index(drop=True)
res_df2 = check_nan(select_df2, link_df)[['order_id','test_na']]
res_df2.to_excel('simp_result2_6may.xlsx', index=False)